from pathlib import Path

from django.core.exceptions import ValidationError


SUPPORTED_EXTENSIONS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".txt": "text",
    ".pptx": "PowerPoint",
}


def parse_uploaded_files(files):
    sections = []
    for uploaded_file in files:
        parsed_text = parse_uploaded_file(uploaded_file)
        if parsed_text:
            sections.append(f"--- {uploaded_file.name} ---\n{parsed_text}")
    return "\n\n".join(sections).strip()


def parse_uploaded_file(uploaded_file):
    extension = Path(uploaded_file.name).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValidationError(f"{uploaded_file.name}: unsupported file type. Supported files: {supported}.")

    try:
        if extension == ".txt":
            return parse_text_file(uploaded_file)
        if extension == ".pdf":
            return parse_pdf_file(uploaded_file)
        if extension == ".docx":
            return parse_docx_file(uploaded_file)
        if extension == ".pptx":
            return parse_pptx_file(uploaded_file)
    except ValidationError:
        raise
    except Exception as exc:
        raise ValidationError(f"{uploaded_file.name}: could not parse file ({exc}).") from exc

    return ""


def parse_text_file(uploaded_file):
    data = uploaded_file.read()
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    raise ValidationError(f"{uploaded_file.name}: could not decode text file.")


def parse_pdf_file(uploaded_file):
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ValidationError("PDF parsing requires the pypdf package. Install requirements first.") from exc

    reader = PdfReader(uploaded_file)
    return "\n\n".join((page.extract_text() or "").strip() for page in reader.pages).strip()


def parse_docx_file(uploaded_file):
    try:
        from docx import Document
    except ImportError as exc:
        raise ValidationError("Word parsing requires the python-docx package. Install requirements first.") from exc

    document = Document(uploaded_file)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    table_cells = []
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                table_cells.append(" | ".join(values))
    return "\n".join([*paragraphs, *table_cells]).strip()


def parse_pptx_file(uploaded_file):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValidationError("PowerPoint parsing requires the python-pptx package. Install requirements first.") from exc

    presentation = Presentation(uploaded_file)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        if text_parts:
            slides.append(f"Slide {index}\n" + "\n".join(text_parts))
    return "\n\n".join(slides).strip()
